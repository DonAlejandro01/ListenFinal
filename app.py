from flask import Flask, render_template, request, redirect, url_for, flash
import os
import io
from PIL import Image
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF para extraer texto de PDFs
import openai
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from google.cloud import vision
import logging

logging.basicConfig(level=logging.DEBUG)

# Cargar las variables de entorno desde el archivo .env
load_dotenv()

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads/'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'ppt', 'pptx'}
app.secret_key = os.urandom(24)  # Genera una clave secreta aleatoria

openai.api_key = os.getenv('OPENAI_API_KEY')  # Clave API de OpenAI desde una variable de entorno

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

# Asegurarse de que la carpeta de subidas exista
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

@app.route('/')
def upload_file():
    return render_template('upload.html')

@app.route('/uploader', methods=['GET', 'POST'])
def uploader_file():
    if request.method == 'POST':
        if 'rubric' not in request.files and 'presentation' not in request.files:
            flash('No se encontró la parte del archivo')
            return redirect(request.url)

        rubric_file = request.files['rubric'] if 'rubric' in request.files else None
        presentation_file = request.files['presentation']
        user_type = request.form['user_type']
        presentation_theme = request.form['presentation_theme']
        presentation_type = request.form['presentation_type']

        if presentation_file.filename == '':
            flash('No se seleccionó ningún archivo de presentación')
            return redirect(request.url)

        if user_type not in ['salesperson', 'developer', 'marketing', 'public_speaker', 'elementary', 'secondary', 'university_students'] and (not rubric_file or rubric_file.filename == ''):
            flash('No se seleccionó ningún archivo de rúbrica')
            return redirect(request.url)

        if presentation_file and allowed_file(presentation_file.filename):
            presentation_filename = secure_filename(presentation_file.filename)
            presentation_path = os.path.join(app.config['UPLOAD_FOLDER'], presentation_filename)
            presentation_file.save(presentation_path)

            presentation_text = extract_text_from_ppt(presentation_path)
            inappropriate_content = check_for_inappropriate_content(presentation_text)

            titles = extract_titles(presentation_path)
            subtitles = extract_subtitles(presentation_path)
            body_texts = extract_body_texts(presentation_path, titles, subtitles)
            images = extract_images(presentation_path)
            analyzed_images = [analyze_image_google_cloud(slide_idx, image) for slide_idx, image in images]

            if user_type in ['salesperson', 'developer', 'marketing', 'public_speaker']:
                general_feedback = generate_general_feedback(presentation_theme, presentation_type, titles, subtitles, body_texts, analyzed_images, user_type)
                return render_template('result.html', presentation_score="N/A", specific_feedback=[], general_feedback=general_feedback.split('\n'), inappropriate_content_feedback=inappropriate_content, rubric_table_html=None, user_type=user_type)
            elif user_type in ['elementary', 'secondary', 'university_students']:
                if rubric_file and allowed_file(rubric_file.filename):
                    rubric_filename = secure_filename(rubric_file.filename)
                    rubric_path = os.path.join(app.config['UPLOAD_FOLDER'], rubric_filename)
                    rubric_file.save(rubric_path)

                    # Extraer tabla del PDF de la rúbrica
                    rubric_table = extract_rubric_table(rubric_path)
                    rubric_table_html = table_to_html(rubric_table)

                    rubric_text = extract_text_from_pdf(rubric_path)
                    is_rubric, _ = check_if_rubric(rubric_text)
                
                    if not is_rubric:
                        flash('El archivo cargado no es una rúbrica válida')
                        return redirect(request.url)

                    measures = get_measures(rubric_text)
                    points_type = get_points_type(rubric_text)
                    max_score = calculate_total_score(measures, points_type)  # Calcular el puntaje máximo

                    total_score, specific_feedback = evaluate_presentation(presentation_text, measures, points_type, titles, subtitles, body_texts, analyzed_images, user_type)
                    general_feedback = generate_general_feedback(presentation_theme, presentation_type, titles, subtitles, body_texts, analyzed_images, user_type)

                    grade = convert_score_to_grade(total_score, max_score)  # Convertir el puntaje total en una nota

                    if grade == 7:
                        general_feedback += "\nFelicidades, has obtenido una nota perfecta. ¡Excelente trabajo!"

                    return render_template('result.html', presentation_score=grade, specific_feedback=specific_feedback, general_feedback=general_feedback.split('\n'), inappropriate_content_feedback=inappropriate_content, rubric_table_html=rubric_table_html, user_type=user_type)
                else:
                    flash('No se seleccionó ningún archivo de rúbrica')
                    return redirect(request.url)
            else:
                flash('Tipo de usuario no permitido')
                return redirect(request.url)
        
        flash('Archivo no permitido')
        return redirect(request.url)
    else:
        return render_template('upload.html')



def extract_text_from_pdf(filepath):
    doc = fitz.open(filepath)
    text = ""
    for page in doc:
        text += page.get_text()
    print(text)  # Imprimir el texto extraído en la consola
    return text

def extract_text_from_ppt(filepath):
    prs = Presentation(filepath)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def check_for_inappropriate_content(text):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a content analysis tool. Your job is to detect any false information or inappropriate words in the provided text and suggest better alternatives if needed."},
                {"role": "user", "content": f"Analiza el siguiente texto y detecta cualquier información falsa o palabras inadecuadas. Proporciona las palabras detectadas y sus mejores alternativas:\n\n{text[:3000]}"}
            ]
        )
        analysis_result = response['choices'][0]['message']['content']
        detected_issues = []
        if "false" in analysis_result.lower() or "inappropriate" in analysis_result.lower():
            detected_issues.append(analysis_result)
        return detected_issues
    except Exception as e:
        print(f"Error al analizar el contenido inapropiado: {e}")
        return ["Error al analizar el contenido inapropiado."]

def check_if_rubric(text):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a file analysis tool. Your job is to determine whether the provided text is a rubric used for evaluation or not."},
                {"role": "user", "content": f"Evalua si este texto está relacionado o no a una rúbrica: {text[:3000]}"}
            ]
        )
        is_rubric = "Sí" in response['choices'][0]['message']['content'] or "yes" in response['choices'][0]['message']['content'].lower()
        return is_rubric, None  # Retornar None como segundo valor
    except Exception as e:
        print(f"Error al evaluar la rúbrica: {e}")
        return False, None  # Asegurarse de retornar dos valores

def get_measures(text):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a file analysis tool. Your job is to extract the rubric statements from the provided text. Each statement describes an aspect of the presentation that is being evaluated."},
                {"role": "user", "content": f"Extrae los enunciados de la rúbrica del siguiente texto. Solo proporciona los títulos de las categorías evaluadas, sin frase introductoria, sin frase de conclusión, sin descripción, ni valores numéricos, ni caracteres especiales. Proporciona cada título en una nueva línea: {text[:3000]}"}
            ]
        )
        measures = response['choices'][0]['message']['content']
        measures_list = [measure.strip() for measure in measures.split('\n') if measure.strip()]
        return measures_list
    except Exception as e:
        print(f"Error al extraer las medidas: {e}")
        return []

def get_points_type(text):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a file analysis tool. Determine the type of points used in the provided rubric text and list them in descending order."},
                {"role": "user", "content": f"Extrae el tipo de puntos que se están utilizando en esta rúbrica. No me des texto introductorio, ni de conclusión, ni tampoco descripción. Solo proporciona los valores numéricos en orden descendente: {text[:3000]}"}
            ]
        )
        points_content = response['choices'][0]['message']['content'].strip()
        points_list = [int(point.strip()) for point in points_content.split(',')]
        return sorted(points_list, reverse=True)
    except Exception as e:
        print(f"Error al determinar el tipo de puntos: {e}")
        return []

def calculate_total_score(measures, points_type):
    if not measures or not points_type:
        return 0
    max_point = max(points_type)
    total_score = len(measures) * max_point
    return total_score

def extract_titles(filepath):
    prs = Presentation(filepath)
    titles = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                title = shape.text_frame.text
                if title:
                    titles.append(title)
                    break
    return titles

def extract_subtitles(filepath):
    prs = Presentation(filepath)
    subtitles = []
    for slide in prs.slides:
        subtitle_found = False
        for shape in slide.shapes:
            if shape.has_text_frame and not subtitle_found:
                paragraphs = shape.text_frame.paragraphs
                if len(paragraphs) > 1:
                    subtitles.append(paragraphs[1].text)
                    subtitle_found = True
        if not subtitle_found:
            subtitles.append("")
    return subtitles

def extract_body_texts(filepath, titles, subtitles):
    prs = Presentation(filepath)
    body_texts = []
    for slide_idx, slide in enumerate(prs.slides):
        body_text = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font_size = run.font.size
                    text = run.text.strip()
                    if font_size and font_size < 2400000 and text not in titles and text not in subtitles:
                        body_text.append(text)
        body_texts.append('\n'.join(body_text))
    return body_texts

def extract_images(filepath):
    prs = Presentation(filepath)
    images = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = io.BytesIO(image.blob)
                pil_image = Image.open(image_bytes)
                images.append((slide_idx, pil_image))  # Incluir el número de diapositiva y la imagen

    return images

def analyze_image_google_cloud(slide_idx, image):
    client = vision.ImageAnnotatorClient()
    img_byte_arr = io.BytesIO()
    image.save(img_byte_arr, format='PNG')
    content = img_byte_arr.getvalue()

    image = vision.Image(content=content)
    response = client.label_detection(image=image)
    labels = response.label_annotations

    analyzed_image_info = ', '.join([label.description for label in labels])
    print(f"Slide {slide_idx}: {analyzed_image_info}")  # Incluir el número de diapositiva en la impresión
    
    return slide_idx, analyzed_image_info  # Devolver el número de diapositiva junto con la información analizada

def evaluate_presentation(presentation_text, measures, points_type, titles, subtitles, body_texts, analyzed_images, user_type):
    try:
        measures_str = "\n".join(measures)
        points_str = ", ".join(map(str, points_type))
        titles_str = "\n".join(titles)
        subtitles_str = "\n".join(subtitles)
        body_texts_str = "\n".join(body_texts)
        images_info_str = "\n".join([f"Slide {idx}: {info}" for idx, info in analyzed_images])

        prompt = f"""
        Evalúa la siguiente presentación basada en las medidas y tipos de puntos proporcionados. Ademas debes considerar que el texto proporcionado es apropiado para {user_type} Proporciona un puntaje para cada medida y un feedback específico:

        Tipo de usuario: {user_type}

        Medidas:
        {measures_str}

        Tipos de puntos:
        {points_str}

        Títulos:
        {titles_str}

        Subtítulos:
        {subtitles_str}

        Textos del cuerpo:
        {body_texts_str}

        Información de las imágenes:
        {images_info_str}
        """

        print(f"Prompt enviado a OpenAI: {prompt}")
        
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an evaluation tool. Your job is to evaluate the provided presentation text based on the given rubric measures and point types for the specified user type, and provide scores for each measure and specific feedback."},
                {"role": "user", "content": prompt}
            ]
        )

        response_content = response['choices'][0]['message']['content'].strip()
        print(f"Respuesta de OpenAI: {response_content}")

        lines = response_content.split("\n")
        scores = {}
        feedback_lines = []
        for line in lines:
            if ":" in line:
                measure, score = line.split(":")
                measure = measure.strip()
                try:
                    score = int(score.strip().split("/")[0])
                    scores[measure] = score
                    feedback_lines.append(f"{measure}: {score}")
                except ValueError:
                    scores[measure] = 0
                    feedback_lines.append(f"{measure}: 0")
            else:
                feedback_lines.append(line)

        feedback_list = [item.strip() for item in feedback_lines if item.strip()]

        total_score = sum(scores.values())

        return total_score, feedback_list
    except Exception as e:
        print(f"Error al evaluar la presentación: {e}")
        return 0, ["Error en la evaluación de la presentación."]

def generate_general_feedback(theme, p_type, titles, subtitles, body_texts, images_info, user_type):
    prompt = f"""
    Tema de la presentación: {theme}
    Tipo de presentación: {p_type}
    Tipo de usuario: {user_type}
    
    Títulos: {', '.join(titles)}
    Subtítulos: {', '.join(subtitles)}
    Textos del cuerpo: {', '.join(body_texts)}
    Información de las imágenes: {', '.join(info for _, info in images_info)}

    Basado en la información anterior, proporciona recomendaciones generales para mejorar la presentación y asegurar que sea excelente. Asegurate que el la informacion a comunicar sea acorde para {user_type}
    """

    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are an evaluation assistant. Your job is to provide general feedback to improve the presentation based on the provided details and the specified user type."},
            {"role": "user", "content": prompt}
        ]
    )

    feedback = response['choices'][0]['message']['content'].strip()
    return feedback


def convert_score_to_grade(total_score, max_score):
    if max_score == 0:
        return 0  # Evitar división por cero
    percentage = total_score / max_score
    if percentage >= 0.9:
        return 7
    elif percentage >= 0.8:
        return 6
    elif percentage >= 0.7:
        return 5
    elif percentage >= 0.6:
        return 4
    elif percentage >= 0.5:
        return 3
    elif percentage >= 0.4:
        return 2
    else:
        return 1

import pdfplumber

def extract_rubric_table(pdf_path):
    """
    Extrae una tabla de rúbrica de un archivo PDF, incluyendo tablas que se extienden a lo largo de varias páginas.
    Args:
        pdf_path (str): La ruta al archivo PDF.
    Returns:
        list: Lista de listas que representa la tabla extraída del PDF.
    """
    table = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                extracted_table = page.extract_table()
                if extracted_table:
                    if not table:
                        table.extend(extracted_table)  # Añadir la primera tabla completa
                    else:
                        table.extend(extracted_table[1:])  # Añadir la tabla sin el encabezado
    except Exception as e:
        print(f"Error al extraer la tabla del PDF: {e}")
    
    return table

def table_to_html(table):
    """
    Construye una tabla HTML a partir de la tabla extraída del PDF.
    Args:
        table (list): Lista de listas que representa la tabla extraída del PDF.
    Returns:
        str: Tabla HTML.
    """
    if not table:
        return "<p>No se encontró ninguna tabla en el PDF.</p>"

    header = table[0]
    html_table = '<table class="table table-bordered">'
    
    # Construir el encabezado
    html_table += '<thead><tr>'
    for cell in header:
        html_table += f'<th>{cell}</th>'
    html_table += '</tr></thead>'
    
    # Construir el cuerpo de la tabla
    html_table += '<tbody>'
    for row in table[1:]:
        html_table += '<tr>'
        for cell in row:
            html_table += f'<td>{cell}</td>'
        html_table += '</tr>'
    html_table += '</tbody>'
    
    html_table += '</table>'
    return html_table

if __name__ == '__main__':
    app.run(debug=True)
